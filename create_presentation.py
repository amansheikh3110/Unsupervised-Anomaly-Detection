"""
Generate PowerPoint presentation from PRESENTATION_CONTENT.txt
Run: python create_presentation.py
Output: Anomaly_Detection_Presentation.pptx
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt

# Slide layout indices: 0 = title only, 1 = title + content, 6 = blank
TITLE_ONLY = 0
TITLE_AND_BODY = 1


def add_slide_title_only(prs, title_text):
    """Add a slide with only a title (e.g. section divider)."""
    layout = prs.slide_layouts[TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    return slide


def add_slide_bullets(prs, title_text, bullets, sub_bullets=None):
    """
    Add a slide with title and bullet points.
    bullets: list of strings (main bullets).
    sub_bullets: optional list of lists; sub_bullets[i] = list of strings under bullets[i].
    """
    layout = prs.slide_layouts[TITLE_AND_BODY]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.space_after = Pt(6)
        if sub_bullets is not None and i < len(sub_bullets) and sub_bullets[i]:
            for sub in sub_bullets[i]:
                p2 = tf.add_paragraph()
                p2.text = sub
                p2.level = 1
                p2.space_after = Pt(2)
    return slide


def add_slide_paragraphs(prs, title_text, lines):
    """Add a slide with title and a list of lines (each can be a bullet or block)."""
    layout = prs.slide_layouts[TITLE_AND_BODY]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        # Detect sub-bullet: starts with "–" or "  "
        stripped = line.strip()
        if stripped.startswith("–") or stripped.startswith("-"):
            p.text = stripped.lstrip("–-").strip()
            p.level = 1
        else:
            p.text = stripped
            p.level = 0
        p.space_after = Pt(4)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- 1. TITLE SLIDE -----
    layout0 = prs.slide_layouts[TITLE_AND_BODY]
    s0 = prs.slides.add_slide(layout0)
    s0.shapes.title.text = "Unsupervised Anomaly Detection Using Deep Representation Learning (I-JEPA-Based Industry-Grade AI)"
    ph1 = s0.placeholders[1]
    tf1 = ph1.text_frame
    tf1.text = "Shri Ramdeobaba College of Engineering & Management, Nagpur"
    tf1.add_paragraph().text = "Department of Computer Science & Engineering | Semester VI | Section A"
    tf1.add_paragraph().text = "Group A-5: Aman Sheikh, Aadil Pathan, Aryan Jaiswal"
    tf1.add_paragraph().text = "Guide: Dr. P. Sonsare"

    # ----- 2. INTRODUCTION -----
    add_slide_bullets(prs, "Introduction — Problem & Solution",
        [
            "Problem we are solving:",
            "In manufacturing, visual inspection for defects is critical but costly and error-prone when done manually.",
            "Supervised methods need large labelled defect data (scarce in industry); defects are rare, diverse, and often unseen.",
            "Existing unsupervised (reconstruction-based) methods detect via reconstruction error but often miss high-level semantic defects (structural anomalies, missing components).",
            "High-level idea of our solution:",
            "We build an unsupervised system that learns only from normal (good) images.",
            "Phase 1 (current): Baseline convolutional autoencoder; anomaly score = reconstruction error (MSE).",
            "Phase 2 (planned): I-JEPA learns semantic representations in latent space; anomalies = deviations from learned normal distribution.",
        ],
        sub_bullets=[
            [],
            [], [], [],
            [],
            [], [], [],
        ]
    )

    # ----- 3. MOTIVATION -----
    add_slide_bullets(prs, "Motivation",
        [
            "Why this project is necessary:",
            "Manual inspection: slow, expensive, inconsistent; not scalable.",
            "Supervised defect detection: requires many labelled defect samples; poor generalization to novel defect types (Ruff et al., 2021; Bergmann et al., 2019).",
            "Reconstruction-based methods (e.g. autoencoders): optimize pixel-level reconstruction, not semantic consistency; can fail on structural/semantic anomalies (cite Unreliability of autoencoders paper).",
            "Research gap: Need for methods that learn semantic representations of 'normal' and detect anomalies as deviations in representation space, without defect labels.",
            "Real-world impact: Automated visual inspection, quality control, surface defect detection, semiconductor and packaging verification (Bergmann et al., 2019).",
        ]
    )

    # ----- 4. LITERATURE — Traditional Autoencoder -----
    add_slide_bullets(prs, "Literature Review — Traditional Autoencoder",
        [
            "Core idea: Encoder maps input x to latent z; decoder maps z to reconstruction x̂. Trained to minimize reconstruction loss (MSE). In anomaly detection: train only on normal data; high reconstruction error = anomaly score.",
            "Key equations:",
            "Reconstruction loss (MSE):  L_rec = (1/n) Σ ‖x − x̂‖²",
            "Anomaly score:  s(x) = ‖x − Dec(Enc(x))‖²",
            "Advantages: Unsupervised; no defect labels; simple to implement; interpretable.",
            "Limitations: Pixel-level reconstruction, not semantic fidelity; can reconstruct some anomalies well and some normals poorly; weak on structural/semantic anomalies.",
            "Diagram: Insert figure from traditional_autoencoder.pdf (bottleneck architecture). Cite the paper.",
        ]
    )

    # ----- 5. LITERATURE — Unreliability -----
    add_slide_bullets(prs, "Literature Review — Unreliability of Autoencoders",
        [
            "Core idea: Reconstruction error alone is an unreliable indicator of anomaly when the model is trained only on normal data (statistical and empirical evidence).",
            "Key points: Cases where autoencoders assign low reconstruction error to anomalous samples or high error to normal samples; overlap between normal/anomaly score distributions.",
            "This motivates moving to representation-based methods (latent space, semantic features) rather than pixel reconstruction — e.g. I-JEPA.",
            "Diagram: Insert figure/table from Unreliability of autoencoders.pdf (failure cases or score distributions). Cite the paper.",
        ]
    )

    # ----- 6. LITERATURE — I-JEPA -----
    add_slide_bullets(prs, "Literature Review — I-JEPA (Assran et al., 2023)",
        [
            "Core idea: I-JEPA predicts representations of target (masked) regions from context (visible) regions in a joint embedding space. Does not reconstruct pixels.",
            "Components: Context encoder (visible patches), target encoder (masked patches, EMA-updated), predictor (context → predicted target representation).",
            "Key equations:",
            "z_s = f_θ(s)  (context);  z_t = g_ξ(t)  (target);  Loss: L = ‖p_θ(z_s) − stop_grad(z_t)‖²",
            "Advantages: Semantic, high-level representations; robust to pixel variations; well-suited for anomaly detection in latent space.",
            "Limitations: More complex; requires careful masking and architecture design.",
            "Diagram: Insert main method figure from I-JEPA.pdf (context encoder, target encoder, predictor). Cite Assran et al. (2023).",
        ]
    )

    # ----- 7. LITERATURE — Summary Table -----
    add_slide_bullets(prs, "Literature Review — Comparison",
        [
            "Method: Autoencoder  |  I-JEPA",
            "Learning signal: Pixel reconstruction  |  Latent prediction",
            "Anomaly cue: Reconstruction error  |  Deviation in latent space",
            "Semantic awareness: Limited  |  Strong (Assran et al., 2023)",
        ]
    )

    # ----- 8. PROPOSED METHODOLOGY -----
    add_slide_bullets(prs, "Proposed Methodology — Overview",
        [
            "1. Data: MVTec AD (Bergmann et al., 2019). Train: only normal (good) images per category. Test: normal + defective.",
            "2. Preprocessing: Resize 224×224; ImageNet normalize (mean [0.485, 0.456, 0.406], std [0.229, 0.224, 0.225]).",
            "3. Baseline (current): Train convolutional autoencoder on normal only. Anomaly score = MSE(input, reconstruction). Threshold from validation (best F1).",
            "4. Future (I-JEPA): Train I-JEPA on normal images; extract context encoder embeddings; fit normality model (Mahalanobis, k-NN, or OCSVM). Test: embedding → distance to normal = anomaly score.",
            "5. Evaluation: ROC-AUC, average precision, F1 (labels used only for evaluation).",
        ]
    )

    # ----- 9. FLOW DIAGRAM (text description) -----
    add_slide_bullets(prs, "Proposed Methodology — Flow",
        [
            "Data → MVTec AD (train: good only; test: good + defect)",
            "Preprocessing → Resize 224×224, ImageNet normalize",
            "Model → Current: Autoencoder (Encoder → Latent z → Decoder → x̂). Future: I-JEPA (Patchify → Mask → Context/Target encoders → Predictor → latent MSE).",
            "Training → Normal images only. Loss: MSE (AE) or latent MSE (I-JEPA). Adam/AdamW.",
            "Inference → AE: score = MSE(x, x̂). I-JEPA: score = distance to normal model (e.g. Mahalanobis).",
            "Evaluation → ROC-AUC, AP, F1; optional heatmap for localization.",
        ]
    )

    # ----- 10. TECH STACK -----
    add_slide_bullets(prs, "Tech Stack",
        [
            "Language & runtime: Python 3.10+",
            "Deep learning: PyTorch 2.7 (CUDA 11.8), torchvision, torchaudio",
            "Data & numerical: NumPy, Pandas, SciPy",
            "ML & evaluation: Scikit-learn (metrics, k-NN, OCSVM)",
            "Visualization & imaging: Matplotlib, Seaborn, OpenCV, Pillow",
            "Web demo: Flask, Werkzeug",
            "Development: tqdm, Jupyter (optional). All dependencies in requirements.txt; reproducible via venv.",
        ]
    )

    # ----- 11. IMPLEMENTATION STATUS -----
    add_slide_bullets(prs, "Current Implementation Status",
        [
            "Phase 0: Environment (Python, PyTorch, CUDA, venv); all libraries verified.",
            "Phase 1: MVTec AD loaded and preprocessed; dataset API (train/test, category, patchify); data exploration and visualizations.",
            "Phase 2: Baseline convolutional autoencoder (224×224 → encoder → latent 256-d → decoder); training pipeline (normal only, Adam, CosineAnnealingLR); evaluation (MSE as anomaly score; ROC-AUC, AP, F1); heatmap generation; single-image check script and Web UI.",
            "Key modules: src/datasets.py, model_autoencoder.py, train_baseline.py, anomaly_eval.py, heatmap_utils.py, backend.py; run_baseline.py, check_image.py, app.py.",
        ]
    )

    # ----- 12. DEMO SUGGESTIONS -----
    add_slide_bullets(prs, "Demo Suggestions",
        [
            "Terminal: python run_baseline.py --category leather --epochs 2; or python check_image.py data/leather/test/cut/002.png --category leather; show score, threshold, Normal/Anomaly, heatmap.",
            "Web UI: python app.py → http://127.0.0.1:5000. Dataset tab: browse and click image to run detection. Tester tab: drag-drop image, Run detection, show result (score, ROC-AUC, F1), heatmap.",
            "Code: Show Autoencoder in model_autoencoder.py and training loop in train_baseline.py (only normal data).",
        ]
    )

    # ----- 13. DATASET -----
    add_slide_bullets(prs, "Dataset — MVTec AD",
        [
            "Citation: Bergmann et al. (2019). The MVTec Anomaly Detection Dataset: A Comprehensive Real-World Dataset for Unsupervised Anomaly Detection. IJCV (use exact from your dataset PDF).",
            "Description: 15 object/texture categories (bottle, cable, capsule, carpet, grid, hazelnut, leather, metal_nut, pill, screw, tile, toothbrush, transistor, wood, zipper). Train: defect-free only; test: good + defective with pixel-level ground-truth masks.",
            "Size: 3,629 training (normal only), 1,725 test. Example leather: 245 train good; 124 test (32 good, 92 defective: color, cut, fold, glue, poke).",
            "Preprocessing: Resize 224×224; ImageNet normalization; optional flip for training. Train on train/good only (unsupervised).",
        ]
    )

    # ----- 14. FUTURE WORK -----
    add_slide_bullets(prs, "Work to Be Done Before Semester End",
        [
            "Foundation: I-JEPA architecture (Assran et al., 2023) — predict masked region representations from context.",
            "1. Implement I-JEPA: context encoder, target encoder (EMA), predictor, semantic-scale masking (Assran et al., 2023).",
            "2. Integrate with existing pipeline (MVTec AD, 224×224, patchify); train on normal only.",
            "3. Extract context encoder embeddings; fit normality models (Mahalanobis, k-NN, OCSVM).",
            "4. Anomaly score = distance to normal model; evaluate ROC-AUC, AP, F1 per category.",
            "5. Comparative experiments: baseline autoencoder vs. I-JEPA on MVTec AD; discuss why I-JEPA expected to outperform (semantic vs. pixel; cite unreliability paper).",
            "6. Optional: ablations (mask ratio, normality model); t-SNE/PCA and heatmaps.",
        ]
    )

    # ----- 15. REFERENCES -----
    add_slide_bullets(prs, "References",
        [
            "Assran, M., Duval, Q., Misra, I., Bojanowski, P., Vincent, P., Rabbat, M., & LeCun, Y. (2023). Self-Supervised Learning from Images with a Joint-Embedding Predictive Architecture. CVPR 2023. (Replace with exact from I-JEPA PDF.)",
            "Bergmann, P., Batzner, K., Fauser, M., Sattlegger, D., & Steger, C. (2019). The MVTec Anomaly Detection Dataset: A Comprehensive Real-World Dataset for Unsupervised Anomaly Detection. International Journal of Computer Vision, 129(4), 1038–1059. (Replace with exact from dataset PDF.)",
            "Traditional autoencoder: [Insert full reference from traditional_autoencoder.pdf.]",
            "Unreliability of autoencoders: [Insert full reference from Unreliability of autoencoders.pdf.]",
        ]
    )

    # ----- 16. THANK YOU -----
    add_slide_title_only(prs, "Thank You\nQuestions?")

    out_path = "Anomaly_Detection_Presentation.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
