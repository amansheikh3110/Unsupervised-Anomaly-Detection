"""
build_ppt.py — Final Presentation PPT Builder
Run: python build_ppt.py
Output: Final_Presentation_GroupA5.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette ────────────────────────────────────────────────────────────
BLACK   = RGBColor(0x1A, 0x1A, 0x2E)   # near-black for headings
DARK    = RGBColor(0x16, 0x21, 0x3E)   # dark navy accent
ACCENT  = RGBColor(0x0F, 0x3D, 0x91)   # strong blue for highlights
TEAL    = RGBColor(0x06, 0x7B, 0xC2)   # medium blue
GREEN   = RGBColor(0x1A, 0x93, 0x6F)   # green for positive results
RED     = RGBColor(0xC0, 0x39, 0x2B)   # red for problems
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)   # orange accent
GRAY    = RGBColor(0x55, 0x55, 0x55)   # body text
LGRAY   = RGBColor(0xAA, 0xAA, 0xAA)   # light gray
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xFF, 0xFF, 0xFF)   # slide background: white

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

BASE = os.path.dirname(os.path.abspath(__file__))

def img(rel):
    p = os.path.join(BASE, rel)
    return p if os.path.exists(p) else None

VIZ = "visualizations"
HM  = "results/heatmaps"
EX  = "results/exploration"

IMGS = {
    "sorted_auroc"   : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.14.21 PM.jpeg"),
    "img_pixel_dual" : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.14.36 PM.jpeg"),
    "auroc_heatmap"  : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.14.46 PM.jpeg"),
    "3way_compare"   : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.14.55 PM.jpeg"),
    "pixel_compare"  : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.15.16 PM.jpeg"),
    "mean_auroc_bar" : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.15.30 PM.jpeg"),
    "scatter"        : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.15.45 PM.jpeg"),
    "improvement"    : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.15.58 PM.jpeg"),
    "boxplot"        : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.16.04 PM.jpeg"),
    "trend"          : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.16.13 PM.jpeg"),
    "pixel_var"      : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.16.23 PM.jpeg"),
    "radar"          : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.16.34 PM.jpeg"),
    "all_models_hm"  : img(f"{VIZ}/WhatsApp Image 2026-04-10 at 10.15.24 PM.jpeg"),
    "dataset_ov"     : img(f"{EX}/dataset_overview.png"),
    "leather_samp"   : img(f"{EX}/leather_samples.png"),
    "patch_viz"      : img(f"{EX}/patchification.png"),
    "hm_cut"         : img(f"{HM}/heatmap_leather_test_cut_007.png"),
    "hm_fold"        : img(f"{HM}/heatmap_leather_test_fold_009.png"),
    "hm_glue"        : img(f"{HM}/heatmap_leather_test_glue_015.png"),
    "hm_good_ae"     : img(f"{HM}/heatmap_leather_test_good_015.png"),
    "hm_ijepa_hz1"   : img(f"{HM}/ijepa_hazelnut_001.png"),
    "hm_ijepa_hz2"   : img(f"{HM}/ijepa_hazelnut_013.png"),
    "hm_ijepa_cb1"   : img(f"{HM}/ijepa_cable_008.png"),
    "hm_ijepa_cb2"   : img(f"{HM}/ijepa_cable_009.png"),
    "hm_ijepa_mn"    : img(f"{HM}/ijepa_metal_nut_013.png"),
    "hm_ijepa_sc"    : img(f"{HM}/ijepa_screw_000.png"),
}

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper functions ──────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)

def bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def txb(slide, text, l, t, w, h, size=18, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def txb_lines(slide, lines, l, t, w, h, size=14, bold=False, color=BLACK,
              align=PP_ALIGN.LEFT, line_space=1.15, first_bold=False):
    """Add multiple lines to one textbox. lines = list of (text, bold, color, size)"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    tf_frame = tf.text_frame
    tf_frame.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, c, s = item, bold, color, size
        elif len(item) == 4:
            txt, b, c, s = item
        elif len(item) == 3:
            txt, b, c = item; s = size
        elif len(item) == 2:
            txt, b = item; c = color; s = size
        else:
            txt = item[0]; b = bold; c = color; s = size

        if first:
            p = tf_frame.paragraphs[0]
            first = False
        else:
            p = tf_frame.add_paragraph()

        p.alignment = align
        if line_space != 1.0:
            from pptx.util import Pt
            from pptx.oxml.ns import qn
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', f'{int(line_space*100000)}')

        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(s)
        run.font.bold  = b
        run.font.color.rgb = c
    return tf

def rect(slide, l, t, w, h, fill_color=ACCENT, alpha=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_img(slide, key, l, t, w, h=None):
    path = IMGS.get(key)
    if not path:
        return
    if h:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def hline(slide, t, color=LGRAY, thickness_pt=0.75):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    line = slide.shapes.add_shape(1, Inches(0.4), Inches(t), Inches(12.5), Inches(0.01))
    line.fill.background()
    line.line.color.rgb = color
    line.line.width = Pt(thickness_pt)

def slide_header(slide, title, subtitle=None, accent_bar=True):
    bg_white(slide)
    if accent_bar:
        r = rect(slide, 0, 0, 13.33, 0.08, ACCENT)
    txb(slide, title, 0.4, 0.18, 12.5, 0.65, size=28, bold=True, color=BLACK)
    if subtitle:
        txb(slide, subtitle, 0.4, 0.82, 12.5, 0.4, size=14, color=GRAY, italic=True)
    if subtitle:
        hline(slide, 1.22)
    else:
        hline(slide, 0.88)

def footer(slide, txt="Group A-5 | Aman Sheikh, Aadil Pathan, Aryan Jaiswal | RCOEM Nagpur | 2025-26"):
    txb(slide, txt, 0.4, 7.25, 12.5, 0.22, size=8, color=LGRAY, align=PP_ALIGN.CENTER)
    rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)

import pptx.enum.shapes

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bg_white(sl)

# Top accent bar
rect(sl, 0, 0, 13.33, 1.1, ACCENT)
txb(sl, "FINAL PRESENTATION", 0, 0.1, 13.33, 0.45, size=13, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Semester VI | Group A-5 | 2025–26", 0, 0.52, 13.33, 0.4, size=11,
    color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

# Main title
txb(sl, "Unsupervised Anomaly Detection", 0.5, 1.3, 12.3, 0.9,
    size=36, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
txb(sl, "Using Deep Representation Learning", 0.5, 2.08, 12.3, 0.8,
    size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txb(sl, "I-JEPA + Pretrained ViT-B/16 | MVTec AD | Real-Time Defect Detection & Localization",
    0.5, 2.88, 12.3, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

hline(sl, 3.55, LGRAY)

# Details
txb(sl, "Shri Ramdeobaba College of Engineering & Management, Nagpur",
    0.5, 3.65, 12.3, 0.38, size=13, color=BLACK, align=PP_ALIGN.CENTER)
txb(sl, "Department of Computer Science & Engineering",
    0.5, 3.98, 12.3, 0.35, size=12, color=GRAY, align=PP_ALIGN.CENTER)

hline(sl, 4.45, LGRAY)

txb_lines(sl, [
    ("Group A-5:   Aman Sheikh  |  Aadil Pathan  |  Aryan Jaiswal", True, BLACK, 13),
    ("Guide:  Dr. P. Sonsare                                        April 2026", False, GRAY, 12),
], 0.5, 4.58, 12.3, 0.85, align=PP_ALIGN.CENTER)

# Result badge
rect(sl, 3.8, 5.6, 5.73, 1.5, RGBColor(0xF0, 0xF7, 0xFF))
txb(sl, "Mean Image AUROC", 3.9, 5.68, 5.5, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER)
txb(sl, "91.9%", 3.9, 5.98, 5.5, 0.7, size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txb(sl, "across all 15 MVTec AD categories", 3.9, 6.6, 5.5, 0.35, size=10, color=GRAY, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — QUICK RECAP: WHERE WE LEFT OFF
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Presentation 1 — Where We Left Off", "Quick recap of what was covered in the first presentation")

# Left panel
rect(sl, 0.4, 1.3, 5.8, 5.8, RGBColor(0xF5, 0xF8, 0xFF))
txb(sl, "What We Presented", 0.6, 1.38, 5.4, 0.42, size=13, bold=True, color=ACCENT)
txb_lines(sl, [
    ("Problem Statement", True, BLACK, 12),
    ("Detecting manufacturing defects without any labelled defect images during training.", False, GRAY, 11),
    ("", False, GRAY, 6),
    ("Dataset: MVTec AD", True, BLACK, 12),
    ("15 industrial categories  |  3,629 normal training images  |  1,725 test images", False, GRAY, 11),
    ("", False, GRAY, 6),
    ("Phase 2 Complete: Autoencoder Baseline", True, BLACK, 12),
    ("• Convolutional Autoencoder (224×224 → 256-dim latent)", False, GRAY, 11),
    ("• Anomaly score = pixel MSE reconstruction error", False, GRAY, 11),
    ("• Tested on: leather category", False, GRAY, 11),
    ("• Result: ROC-AUC = 0.413  ← very poor", True, RED, 11),
    ("", False, GRAY, 6),
    ("Future Plan Announced", True, BLACK, 12),
    ("Replace pixel MSE with I-JEPA — semantic latent-space prediction", False, GRAY, 11),
], 0.6, 1.82, 5.5, 5.2, size=11)

# Right panel
rect(sl, 6.7, 1.3, 6.2, 5.8, RGBColor(0xFF, 0xF5, 0xF5))
txb(sl, "The Gap We Identified", 6.9, 1.38, 5.9, 0.42, size=13, bold=True, color=RED)
txb_lines(sl, [
    ("Why Autoencoder Failed on Leather:", True, BLACK, 12),
    ("• AE learned to reconstruct ALL images well", False, GRAY, 11),
    ("• Including anomalous ones → low MSE → missed", False, GRAY, 11),
    ("• AUC = 0.41: worse than random guessing (0.50)", True, RED, 11),
    ("", False, GRAY, 6),
    ("Root Cause:", True, BLACK, 12),
    ("Pixel-level loss cannot capture semantic meaning.", False, GRAY, 11),
    ("A cut in leather has similar pixel statistics to", False, GRAY, 11),
    ("a textured fold — the AE cannot tell them apart.", False, GRAY, 11),
    ("", False, GRAY, 6),
    ("The Solution We Planned:", True, BLACK, 12),
    ("Learn semantic representations, not pixels.", False, GRAY, 11),
    ("Flag anomalies as distances in feature space,", False, GRAY, 11),
    ("not as pixel reconstruction errors.", False, GRAY, 11),
    ("", False, GRAY, 6),
    ("→  This is exactly what I-JEPA enables.", True, ACCENT, 12),
], 6.9, 1.82, 5.9, 5.2, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROJECT JOURNEY TIMELINE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Our Journey — 4 Phases from Setup to Production-Grade AI")

colors_ph = [RGBColor(0xAA, 0xAA, 0xAA), RGBColor(0xE6, 0x7E, 0x22),
             RGBColor(0x0F, 0x3D, 0x91), RGBColor(0x1A, 0x93, 0x6F)]
phases = [
    ("PHASE 0 & 1", "Feb 2026", "Environment + Dataset",
     "Python, PyTorch, CUDA\nMVTec AD loaded\n15 categories explored\nData pipeline built", "—"),
    ("PHASE 2", "Feb 2026", "Autoencoder Baseline",
     "Conv AE (224→256→224)\nTrained on leather\nPixel MSE anomaly score\nROC-AUC: 0.413", "0.413"),
    ("PHASE 3", "Mar 2026", "Custom I-JEPA Training",
     "ViT-Small/16 built\nBlock masking strategy\nSelf-supervised training\nI-JEPA methodology", "~0.75"),
    ("PHASE 4", "Apr 2026", "I-JEPA + Pretrained ViT",
     "Pretrained ViT-B/16\nPatchCore-style k-NN\nAll 15 categories\nMean AUROC: 0.919", "0.919"),
]
xs = [0.35, 3.45, 6.55, 9.65]
for i, (ph, date, title, desc, auc) in enumerate(phases):
    c = colors_ph[i]
    rect(sl, xs[i], 1.35, 3.25, 5.7, c)
    txb(sl, ph, xs[i]+0.1, 1.42, 3.05, 0.38, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, date, xs[i]+0.1, 1.78, 3.05, 0.28, size=10, color=RGBColor(0xDD,0xDD,0xFF), align=PP_ALIGN.CENTER)
    hline(sl, 2.12)
    txb(sl, title, xs[i]+0.1, 2.18, 3.05, 0.45, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, desc, xs[i]+0.15, 2.68, 2.95, 2.8, size=10, color=RGBColor(0xEE,0xEE,0xFF), align=PP_ALIGN.LEFT)
    # AUC badge
    rect(sl, xs[i]+0.5, 5.5, 2.25, 1.1, RGBColor(0xFF,0xFF,0xFF) if i==3 else RGBColor(0,0,0,))
    if i == 3:
        txb(sl, "Mean AUROC", xs[i]+0.5, 5.52, 2.25, 0.33, size=9, color=ACCENT, align=PP_ALIGN.CENTER)
        txb(sl, auc, xs[i]+0.5, 5.82, 2.25, 0.55, size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    else:
        txb(sl, "AUROC", xs[i]+0.5, 5.52, 2.25, 0.33, size=9, color=RGBColor(0xCC,0xCC,0xFF), align=PP_ALIGN.CENTER)
        txb(sl, auc, xs[i]+0.5, 5.82, 2.25, 0.55, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow labels
for x in [3.25, 6.35, 9.45]:
    txb(sl, "→", x-0.1, 3.8, 0.5, 0.5, size=22, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT IS I-JEPA? CORE CONCEPT
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "I-JEPA — Core Concept", "Image-based Joint-Embedding Predictive Architecture  (Assran et al., CVPR 2023)")

# Left: key idea
rect(sl, 0.4, 1.28, 6.0, 5.85, RGBColor(0xF5, 0xF8, 0xFF))
txb(sl, "The Key Idea", 0.6, 1.36, 5.7, 0.42, size=14, bold=True, color=ACCENT)
txb_lines(sl, [
    ("What does I-JEPA do?", True, BLACK, 12),
    ("Predicts the semantic representations of masked image", False, GRAY, 11),
    ("regions — entirely in latent space, not pixel space.", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("No pixel reconstruction.", True, RED, 12),
    ("The model never compares raw pixels in its loss.", False, GRAY, 11),
    ("It compares embedding vectors only.", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("What this forces the model to learn:", True, BLACK, 12),
    ("• Structure and shape of objects", False, GRAY, 11),
    ("• Texture patterns at a semantic level", False, GRAY, 11),
    ("• What 'normal' looks like in feature space", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("Training objective:", True, BLACK, 12),
    ("Given visible context patches → predict the", False, GRAY, 11),
    ("embedding of masked target patches.", False, GRAY, 11),
    ("Loss:  L = ‖ Predictor(z_ctx) − stop_grad(z_tgt) ‖²", True, ACCENT, 11),
    ("", False, GRAY, 5),
    ("Reference: Assran et al. (2023), CVPR", False, LGRAY, 10),
], 0.6, 1.82, 5.7, 5.2, size=11)

# Right: comparison table
rect(sl, 6.9, 1.28, 6.0, 5.85, RGBColor(0xFF, 0xFF, 0xFF))
txb(sl, "Autoencoder  vs  I-JEPA", 6.9, 1.36, 5.8, 0.42, size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

rows = [
    ("", "Autoencoder", "I-JEPA + ViT"),
    ("Predicts", "Pixels", "Latent representations"),
    ("Loss signal", "Pixel MSE", "Embedding MSE"),
    ("Learns", "Low-level texture", "Semantic structure"),
    ("Anomaly cue", "High recon. error", "High embedding distance"),
    ("Semantic awareness", "Limited", "Strong"),
    ("Mean AUROC (ours)", "0.659", "0.919 ✓"),
]
ys = [1.85, 2.35, 2.82, 3.29, 3.76, 4.23, 4.7]
for i, (a, b, c) in enumerate(rows):
    bg = RGBColor(0xE8, 0xF0, 0xFF) if i == 0 else (RGBColor(0xF8, 0xF8, 0xF8) if i % 2 == 0 else WHITE)
    rect(sl, 6.9, ys[i], 6.0, 0.46, bg)
    bold_r = (i == 0)
    col_a = ACCENT if i == 0 else BLACK
    col_b = RED if (i == len(rows)-1) else (BLACK if i == 0 else GRAY)
    col_c = GREEN if (i == len(rows)-1) else (BLACK if i == 0 else GRAY)
    txb(sl, a, 7.0, ys[i]+0.08, 1.6, 0.35, size=11, bold=bold_r or (i==0), color=col_a)
    txb(sl, b, 8.65, ys[i]+0.08, 2.1, 0.35, size=11, bold=bold_r, color=col_b, align=PP_ALIGN.CENTER)
    txb(sl, c, 10.8, ys[i]+0.08, 2.0, 0.35, size=11, bold=bold_r, color=col_c, align=PP_ALIGN.CENTER)

txb_lines(sl, [
    ("Key insight: Predicting in latent (semantic) space forces", False, GRAY, 10),
    ("the model to understand what normal objects look like,", False, GRAY, 10),
    ("not just memorize pixel patterns.", False, GRAY, 10),
], 6.9, 5.25, 5.9, 0.85, size=10)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — I-JEPA ARCHITECTURE DEEP DIVE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "I-JEPA Architecture — Three Components")

cols = [
    ("Context Encoder\n(ViT-Small/16)", ACCENT, [
        "Input: Visible (unmasked) patches",
        "Architecture: ViT-Small/16",
        "  12 layers, 384-dim, 6 heads",
        "Parameters: 21.7M",
        "Training: Gradient descent",
        "",
        "Output: z_ctx",
        "Semantic representation",
        "of visible context patches",
        "",
        "This is the main encoder",
        "that learns from data.",
    ]),
    ("Target Encoder\n(EMA copy)", TEAL, [
        "Input: ALL patches (including masked)",
        "Architecture: Identical to context encoder",
        "Parameters: 21.7M (shared init)",
        "Training: NO gradient!",
        "",
        "Updated via EMA:",
        "  ξ ← m·ξ + (1-m)·θ",
        "  momentum m: 0.996 → 1.0",
        "",
        "Output: z_tgt",
        "Stable target representations",
        "(stop-gradient applied)",
    ]),
    ("Predictor\n(Narrow ViT)", GREEN, [
        "Input: z_ctx + positional info",
        "       of masked regions",
        "Architecture: Narrow ViT",
        "  6 layers, 192-dim",
        "Parameters: 2.9M",
        "",
        "Task: predict z_tgt from z_ctx",
        "",
        "Loss:",
        "  L = ‖Pred(z_ctx) − z_tgt‖²",
        "",
        "Total trainable: ~24.6M params",
    ]),
]
xs = [0.35, 4.7, 9.05]
for i, (title, c, items) in enumerate(cols):
    rect(sl, xs[i], 1.28, 4.0, 5.85, c)
    txb(sl, title, xs[i]+0.1, 1.35, 3.8, 0.65, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    lines = [(it, False, WHITE if it else WHITE, 10 if it else 5) for it in items]
    txb_lines(sl, lines, xs[i]+0.2, 2.05, 3.7, 4.9, size=10)

# arrows
for x in [4.45, 8.8]:
    txb(sl, "→", x-0.15, 3.9, 0.5, 0.5, size=24, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)

# Bottom note
rect(sl, 0.35, 7.1, 12.6, 0.3, RGBColor(0xEE, 0xF4, 0xFF))
txb(sl, "Image: 224×224 → 196 patches (14×14 grid, 16×16px each)  |  Patch embedding dim: 768  |  No masking at inference — all 196 patches used",
    0.5, 7.12, 12.3, 0.26, size=9, color=DARK, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BLOCK MASKING STRATEGY
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Block Masking Strategy — How I-JEPA Selects Which Patches to Predict")

# Left: explanation
rect(sl, 0.4, 1.28, 5.8, 5.85, RGBColor(0xF5, 0xF8, 0xFF))
txb(sl, "Why Masking Strategy Matters", 0.6, 1.36, 5.5, 0.42, size=13, bold=True, color=ACCENT)
txb_lines(sl, [
    ("Random individual patch masking (like BERT):", True, BLACK, 11),
    ("→ Model fills in local texture patterns", False, GRAY, 11),
    ("→ Low-level learning, not semantic", False, RED, 11),
    ("", False, GRAY, 5),
    ("I-JEPA Block Masking (our approach):", True, BLACK, 11),
    ("→ Mask large contiguous rectangular regions", False, GRAY, 11),
    ("→ Forces model to understand object structure", False, GREEN, 11),
    ("→ Cannot be solved by texture interpolation alone", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("Our Implementation (src/masks.py):", True, BLACK, 12),
    ("  BlockMaskGenerator", False, ACCENT, 11),
    ("", False, GRAY, 4),
    ("• Target blocks: 4 per image", False, GRAY, 11),
    ("• Block scale: 15%–20% of image area", False, GRAY, 11),
    ("• Aspect ratio: 0.75 – 1.5", False, GRAY, 11),
    ("• Context: ~55–60% of patches visible", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("Key insight:", True, BLACK, 11),
    ("Semantic-scale masking is what forces I-JEPA to", False, GRAY, 11),
    ("learn high-level structure, not textures.", False, GRAY, 11),
], 0.6, 1.82, 5.6, 5.25, size=11)

# Right: patch grid visualization
rect(sl, 6.7, 1.28, 6.2, 5.85, WHITE)
txb(sl, "14×14 Patch Grid (196 patches per image)", 6.9, 1.36, 5.9, 0.42, size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Draw 14x14 grid manually
cell_w = 0.32
cell_h = 0.30
grid_l = 7.1
grid_t = 1.88

# Define 4 masked blocks (in grid coordinates)
masked_blocks = [
    (1, 1, 4, 3),   # (col, row, width, height)
    (6, 2, 4, 3),
    (2, 6, 3, 4),
    (8, 7, 4, 4),
]
masked_cells = set()
for bc, br, bw, bh in masked_blocks:
    for r in range(br, br+bh):
        for c in range(bc, bc+bw):
            masked_cells.add((r, c))

for row in range(14):
    for col in range(14):
        cl = grid_l + col * cell_w
        ct = grid_t + row * cell_h
        is_masked = (row, col) in masked_cells
        color = RGBColor(0xC0, 0x39, 0x2B) if is_masked else RGBColor(0xD5, 0xE8, 0xF5)
        shape = sl.shapes.add_shape(1, Inches(cl), Inches(ct), Inches(cell_w-0.01), Inches(cell_h-0.01))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = WHITE
        shape.line.width = Pt(0.5)

txb_lines(sl, [
    ("  Red = 4 masked target blocks (to predict)", True, RED, 10),
    ("  Blue = visible context patches (input)", True, TEAL, 10),
    ("", False, GRAY, 4),
    ("The predictor must infer what the red regions", False, GRAY, 10),
    ("look like semantically from the blue context.", False, GRAY, 10),
], 6.9, 6.2, 5.9, 0.95, size=10)

if IMGS.get("patch_viz"):
    pass  # grid drawn manually above

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PHASE 3: CUSTOM I-JEPA TRAINING
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Phase 3 — Custom I-JEPA Training from Scratch on MVTec")

# Top row: 3 info boxes
boxes3 = [
    ("Training Setup", ACCENT, [
        "Data: train/good only (no defects!)",
        "Input: 224×224 RGB → 196 patches",
        "Normalization: ImageNet mean/std",
        "Optimizer: AdamW",
        "Epochs: 100 per category",
        "HW: NVIDIA GTX 1650 (4GB)",
        "Time: ~2–4 hours per category",
    ]),
    ("What Gets Saved", TEAL, [
        "checkpoints/ijepa_small_<cat>.pth",
        "→ Context encoder weights",
        "→ Encodes 'what normal looks like'",
        "   in semantic patch space",
        "",
        "Checkpoint saved when loss",
        "improves (any epoch with Ctrl+C",
        "still gives a valid checkpoint)",
    ]),
    ("Loss Behavior", GREEN, [
        "Loss fluctuates between epochs.",
        "This is EXPECTED and correct.",
        "",
        "Why? Random block masks differ",
        "every epoch → difficulty varies.",
        "",
        "Watch the 20–30 epoch trend,",
        "not individual values.",
        "Converging trend = learning. ✓",
    ]),
]
xs3 = [0.35, 4.62, 8.88]
for i, (title, c, items) in enumerate(boxes3):
    rect(sl, xs3[i], 1.28, 4.0, 4.0, c)
    txb(sl, title, xs3[i]+0.15, 1.36, 3.7, 0.42, size=12, bold=True, color=WHITE)
    lines3 = [(it, False, WHITE, 10) for it in items]
    txb_lines(sl, lines3, xs3[i]+0.2, 1.85, 3.65, 3.3, size=10)

# Bottom: command + key insight
rect(sl, 0.35, 5.42, 7.5, 1.7, RGBColor(0xF5, 0xF5, 0xF5))
txb(sl, "Key Commands", 0.55, 5.5, 7.1, 0.38, size=12, bold=True, color=BLACK)
txb_lines(sl, [
    ("python run_ijepa.py --categories hazelnut --epochs 100", False, ACCENT, 10),
    ("python run_ijepa.py --categories leather hazelnut cable --epochs 100", False, ACCENT, 10),
    ("python run_ijepa.py --categories hazelnut --no_train   # evaluate only", False, GRAY, 10),
], 0.55, 5.92, 7.1, 1.1, size=10)

rect(sl, 8.15, 5.42, 4.85, 1.7, RGBColor(0xFF, 0xF0, 0xE8))
txb(sl, "Key Insight from Phase 3:", 8.35, 5.5, 4.5, 0.38, size=12, bold=True, color=ORANGE)
txb_lines(sl, [
    ("Training a ViT from scratch on only 300–400 normal", False, BLACK, 10),
    ("images produces limited semantic features.", False, BLACK, 10),
    ("→ This led us to Phase 4: Pretrained ViT backbone.", True, ORANGE, 11),
], 8.35, 5.92, 4.55, 1.1, size=10)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — THE KEY INSIGHT: PRETRAINED VIT
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "The Critical Insight — From Scratch Training to Pretrained ViT")

# Left: problem
rect(sl, 0.4, 1.28, 5.8, 5.85, RGBColor(0xFF, 0xF3, 0xF3))
txb(sl, "Problem with Scratch Training", 0.6, 1.36, 5.5, 0.42, size=13, bold=True, color=RED)
txb_lines(sl, [
    ("Why scratch training struggled:", True, BLACK, 12),
    ("", False, GRAY, 5),
    ("• ViT-Small needs large diverse data to develop", False, GRAY, 11),
    ("  rich semantic features (designed for ImageNet scale)", False, GRAY, 11),
    ("", False, GRAY, 4),
    ("• We train on ~300–400 normal images per category", False, GRAY, 11),
    ("  → severely limited data for self-supervised ViT", False, GRAY, 11),
    ("", False, GRAY, 4),
    ("• Result: patch representations are not distinct enough", False, GRAY, 11),
    ("  to clearly separate normal vs anomaly patches", False, GRAY, 11),
    ("  in the k-NN distance space", False, GRAY, 11),
    ("", False, GRAY, 4),
    ("• I-JEPA was originally designed for ImageNet-scale", False, GRAY, 11),
    ("  self-supervised pretraining", False, GRAY, 11),
    ("", False, GRAY, 6),
    ("Mean AUROC (scratch):  ~0.75", True, RED, 13),
], 0.6, 1.82, 5.5, 5.25, size=11)

# Right: solution
rect(sl, 6.7, 1.28, 6.2, 5.85, RGBColor(0xF0, 0xFF, 0xF4))
txb(sl, "The Solution: Pretrained ViT-B/16", 6.9, 1.36, 5.9, 0.42, size=13, bold=True, color=GREEN)
txb_lines(sl, [
    ("Use torchvision's ViT-B/16 pretrained on ImageNet", True, BLACK, 12),
    ("(1.28 million images, 1000 classes)", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("Why this works immediately:", True, BLACK, 12),
    ("• Already encodes shape, texture, structure, edges", False, GRAY, 11),
    ("• Features generalise to industrial images", False, GRAY, 11),
    ("• Normal patches cluster tightly in feature space", False, GRAY, 11),
    ("• Anomaly patches appear as clear outliers", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("This is still I-JEPA's principle:", True, BLACK, 12),
    ("• Semantic patch features (not pixels)", False, GRAY, 11),
    ("• k-NN deviation scoring in feature space", False, GRAY, 11),
    ("• Same interface: extract_features() unchanged", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("Relationship to I-JEPA:", True, BLACK, 12),
    ("I-JEPA philosophy + pretrained ViT backbone", False, GRAY, 11),
    ("= production-grade anomaly detection", False, GRAY, 11),
    ("", False, GRAY, 6),
    ("Mean AUROC (pretrained): 0.919  ✓", True, GREEN, 13),
], 6.9, 1.82, 5.9, 5.25, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FINAL ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Final System Architecture — I-JEPA + Pretrained ViT-B/16")

# Draw pipeline boxes
pipeline = [
    (0.4,  1.35, 2.2, 1.1, "Input Image\n224×224 RGB",               RGBColor(0xDD,0xEE,0xFF), BLACK),
    (2.85, 1.35, 2.5, 1.1, "Patch Embedding\n196 tokens × 768-dim",  RGBColor(0xCC,0xE5,0xFF), BLACK),
    (5.6,  1.35, 2.5, 1.1, "Pretrained ViT-B/16\n12 layers, 86M params", ACCENT, WHITE),
    (8.35, 1.35, 2.5, 1.1, "196 × 768\nPatch Features",              RGBColor(0xCC,0xE5,0xFF), BLACK),
    (11.1, 1.35, 1.8, 1.1, "L2 Norm",                                RGBColor(0xDD,0xEE,0xFF), BLACK),
]
for l, t, w, h, txt, c, tc in pipeline:
    shape = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = c
    shape.line.color.rgb = ACCENT; shape.line.width = Pt(1.0)
    tf = shape.text_frame; tf.word_wrap = True
    from pptx.enum.text import PP_ALIGN as PA
    from pptx.util import Pt as PPt
    p = tf.paragraphs[0]; p.alignment = PA.CENTER
    p.add_run().text = txt
    p.runs[0].font.size = PPt(10)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = tc
for x in [2.7, 5.45, 8.2, 10.95]:
    txb(sl, "→", x, 1.72, 0.35, 0.4, size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

hline(sl, 2.65)

# Two-column bottom: Training path vs Inference path
rect(sl, 0.4, 2.78, 6.0, 4.3, RGBColor(0xF0, 0xF7, 0xFF))
txb(sl, "TRAINING: Normal Images Only", 0.6, 2.86, 5.7, 0.42, size=12, bold=True, color=ACCENT)
txb_lines(sl, [
    ("For each train/good image:", True, BLACK, 11),
    ("  1. Extract 196 patch features via ViT", False, GRAY, 11),
    ("  2. Collect all patches across all normal images", False, GRAY, 11),
    ("  3. Random subsample 25% = Coreset", False, GRAY, 11),
    ("  4. Build k-NN index (k=9) on coreset", False, GRAY, 11),
    ("  5. Save memory bank as .pkl checkpoint", False, GRAY, 11),
    ("", False, GRAY, 4),
    ("Memory bank = compact representation of", True, BLACK, 11),
    ("'what normal patches look like'", False, GRAY, 11),
    ("Saved: checkpoints/ijepa_detector_<cat>.pkl", False, ACCENT, 10),
], 0.6, 3.32, 5.7, 3.6, size=11)

rect(sl, 6.9, 2.78, 6.0, 4.3, RGBColor(0xF0, 0xFF, 0xF4))
txb(sl, "INFERENCE: Test Image Scoring", 6.9, 2.86, 5.8, 0.42, size=12, bold=True, color=GREEN)
txb_lines(sl, [
    ("1. Extract 196 patch features from test image", False, GRAY, 11),
    ("2. For each patch: find k=9 nearest neighbours", False, GRAY, 11),
    ("   in the memory bank", False, GRAY, 11),
    ("3. Per-patch score = distance to nearest normal patch", False, GRAY, 11),
    ("4. Image score = max patch-level score", False, GRAY, 11),
    ("   (one anomalous patch = anomalous image)", False, GRAY, 11),
    ("", False, GRAY, 4),
    ("score = max_i { min_j ||f_i − f_j||₂ }", True, ACCENT, 11),
    ("", False, GRAY, 4),
    ("5. Heatmap: reshape 196 scores → 14×14 grid", False, GRAY, 11),
    ("   → upsample to 224×224 with bilinear interp.", False, GRAY, 11),
    ("Output: score + heatmap PNG", False, GREEN, 11),
], 6.9, 3.32, 5.8, 3.6, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DATASET
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "MVTec Anomaly Detection Dataset", "Bergmann et al. (2019) | 15 Industrial Categories | Real-world manufacturing defects")

# Left: dataset overview image
if IMGS["dataset_ov"]:
    add_img(sl, "dataset_ov", 0.4, 1.28, 7.0, 4.5)
else:
    rect(sl, 0.4, 1.28, 7.0, 4.5, RGBColor(0xEE, 0xEE, 0xEE))
    txb(sl, "results/exploration/dataset_overview.png", 0.6, 3.0, 6.5, 0.5, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Right: stats
rect(sl, 7.8, 1.28, 5.1, 0.55, ACCENT)
txb(sl, "Dataset Statistics", 7.9, 1.34, 5.0, 0.42, size=12, bold=True, color=WHITE)

stats_rows = [
    ("Category", "Train", "Test", "Good", "Defect"),
    ("bottle", "209", "83", "20", "63"),
    ("cable", "224", "150", "58", "92"),
    ("carpet", "280", "117", "28", "89"),
    ("grid", "264", "78", "21", "57"),
    ("hazelnut", "391", "110", "40", "70"),
    ("leather", "245", "124", "32", "92"),
    ("metal_nut", "220", "115", "22", "93"),
    ("pill", "267", "167", "26", "141"),
    ("screw", "320", "160", "41", "119"),
    ("...+5 more", "...", "...", "...", "..."),
    ("TOTAL (15)", "3,629", "1,725", "467", "1,258"),
]
col_xs = [7.82, 9.18, 9.85, 10.52, 11.2]
col_ws = [1.3,  0.62, 0.62, 0.62,  0.62]
for ri, row in enumerate(stats_rows):
    t_y = 1.87 + ri * 0.40
    bg = ACCENT if ri == 0 else (RGBColor(0xF0,0xF5,0xFF) if ri % 2 == 0 else WHITE)
    bg = RGBColor(0x1A,0x93,0x6F) if ri == len(stats_rows)-1 else bg
    rect(sl, 7.8, t_y, 4.98, 0.39, bg)
    for ci, (cx, cw, val) in enumerate(zip(col_xs, col_ws, row)):
        c = WHITE if (ri == 0 or ri == len(stats_rows)-1) else BLACK
        txb(sl, val, cx, t_y+0.06, cw, 0.28, size=9,
            bold=(ri == 0 or ri == len(stats_rows)-1),
            color=c, align=PP_ALIGN.CENTER)

# Bottom key rule
rect(sl, 0.4, 5.92, 12.5, 1.15, RGBColor(0xFF, 0xF3, 0xE0))
txb(sl, "FUNDAMENTAL RULE", 0.6, 5.98, 3.5, 0.35, size=11, bold=True, color=ORANGE)
txb_lines(sl, [
    ("Training uses ONLY train/good images (normal).  Labels are used ONLY for evaluation (ROC-AUC, F1).", True, BLACK, 12),
    ("No defect image is ever shown to the model during training.  This is what makes it truly unsupervised.", False, GRAY, 11),
], 0.6, 6.36, 12.2, 0.65, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — AUTOENCODER RESULTS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Phase 2 Results — Autoencoder Baseline (All 15 Categories)")

# Table
ae_data = [
    ("bottle","0.67","0.60"), ("cable","0.37","0.53"), ("capsule","0.58","0.81"),
    ("carpet","0.43","0.54"), ("grid","0.71","0.55"), ("hazelnut","0.98","0.97"),
    ("leather","0.84","1.00"), ("metal_nut","0.44","0.53"), ("pill","0.64","0.89"),
    ("screw","0.64","0.96"), ("tile","0.73","0.76"), ("toothbrush","0.65","0.81"),
    ("transistor","0.53","0.54"), ("wood","0.93","0.83"), ("zipper","0.72","0.82"),
    ("MEAN","0.659","0.737"),
]
rect(sl, 0.4, 1.28, 5.5, 0.5, ACCENT)
for ci, (label, cx) in enumerate([("Category", 0.42), ("Image AUROC", 2.02), ("Pixel AUROC", 3.82)]):
    txb(sl, label, cx, 1.34, 1.75, 0.38, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, (cat, img_a, pix_a) in enumerate(ae_data):
    ty = 1.82 + ri * 0.33
    is_last = ri == len(ae_data) - 1
    bg = RGBColor(0x1A,0x93,0x6F) if is_last else (RGBColor(0xF5,0xF8,0xFF) if ri%2==0 else WHITE)
    rect(sl, 0.4, ty, 5.5, 0.32, bg)
    tc = WHITE if is_last else BLACK
    val_img = float(img_a)
    val_col = RED if (not is_last and val_img < 0.6) else (GREEN if (not is_last and val_img >= 0.9) else (WHITE if is_last else BLACK))
    txb(sl, cat.replace("_"," "), 0.5, ty+0.05, 1.55, 0.24, size=10, bold=is_last, color=tc)
    txb(sl, img_a, 2.02, ty+0.05, 1.75, 0.24, size=10, bold=is_last, color=val_col if not is_last else WHITE, align=PP_ALIGN.CENTER)
    txb(sl, pix_a, 3.82, ty+0.05, 1.75, 0.24, size=10, bold=is_last, color=tc, align=PP_ALIGN.CENTER)

# Right: observations
rect(sl, 6.3, 1.28, 6.6, 5.85, RGBColor(0xFF, 0xF8, 0xF0))
txb(sl, "Key Observations", 6.5, 1.36, 6.3, 0.42, size=13, bold=True, color=ORANGE)
txb_lines(sl, [
    ("Catastrophic failures:", True, RED, 12),
    ("• Cable: 0.37 — WORSE than random (0.50)", True, RED, 11),
    ("  AE reconstructs defective cables well", False, GRAY, 11),
    ("• Carpet: 0.43 — near random", True, RED, 11),
    ("• Metal nut: 0.44 — near random", True, RED, 11),
    ("", False, GRAY, 5),
    ("Best cases:", True, GREEN, 12),
    ("• Hazelnut: 0.98 — simple texture works OK", False, GRAY, 11),
    ("• Wood: 0.93 — uniform surface helps AE", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("Verdict:", True, BLACK, 12),
    ("• Mean Image AUROC: 0.659", True, RED, 12),
    ("• Extreme variance: 0.37 to 0.98", False, GRAY, 11),
    ("• Completely unreliable for production", True, RED, 11),
    ("", False, GRAY, 5),
    ("Root cause:", True, BLACK, 12),
    ("Pixel MSE cannot capture semantic meaning.", False, GRAY, 11),
    ("Autoencoders generalize too well — they", False, GRAY, 11),
    ("learn to reconstruct everything, including", False, GRAY, 11),
    ("anomalies, especially textured ones.", False, GRAY, 11),
], 6.5, 1.82, 6.3, 5.2, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — I-JEPA + VIT RESULTS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Phase 4 Results — I-JEPA + Pretrained ViT-B/16 (All 15 Categories)")

# Left: AUROC heatmap image
if IMGS["auroc_heatmap"]:
    add_img(sl, "auroc_heatmap", 0.35, 1.28, 6.5, 5.2)
else:
    rect(sl, 0.35, 1.28, 6.5, 5.2, RGBColor(0xEE,0xEE,0xEE))
    txb(sl, "AUROC Heatmap (Image vs Pixel)", 0.5, 3.5, 6.2, 0.5, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Right: table
vit_data = [
    ("bottle","0.999","0.905"), ("cable","0.940","0.886"), ("capsule","0.872","0.903"),
    ("carpet","0.970","0.930"), ("grid","0.903","0.803"), ("hazelnut","0.984","0.955"),
    ("leather","1.000","0.924"), ("metal_nut","0.987","0.960"), ("pill","0.913","0.939"),
    ("screw","0.717","0.815"), ("tile","0.991","0.883"), ("toothbrush","0.992","0.909"),
    ("transistor","0.736","0.928"), ("wood","0.980","0.835"), ("zipper","0.906","0.805"),
    ("MEAN","0.919","0.892"),
]
rect(sl, 7.1, 1.28, 5.85, 0.5, ACCENT)
for ci, (label, cx) in enumerate([("Category", 7.12), ("Image AUROC", 8.65), ("Pixel AUROC", 10.45)]):
    txb(sl, label, cx, 1.34, 1.75, 0.38, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, (cat, img_a, pix_a) in enumerate(vit_data):
    ty = 1.82 + ri * 0.33
    is_last = ri == len(vit_data) - 1
    bg = RGBColor(0x1A,0x93,0x6F) if is_last else (RGBColor(0xF0,0xFF,0xF4) if ri%2==0 else WHITE)
    rect(sl, 7.1, ty, 5.85, 0.32, bg)
    tc = WHITE if is_last else BLACK
    val_img = float(img_a)
    val_col = ORANGE if (not is_last and val_img < 0.80) else (GREEN if not is_last else WHITE)
    txb(sl, cat.replace("_"," "), 7.2, ty+0.05, 1.4, 0.24, size=10, bold=is_last, color=tc)
    txb(sl, img_a, 8.65, ty+0.05, 1.75, 0.24, size=10, bold=(is_last or val_img>=0.98),
        color=val_col if not is_last else WHITE, align=PP_ALIGN.CENTER)
    txb(sl, pix_a, 10.45, ty+0.05, 1.75, 0.24, size=10, bold=is_last, color=tc, align=PP_ALIGN.CENTER)

# Callout
rect(sl, 7.1, 7.0, 5.85, 0.38, RGBColor(0xE8, 0xF8, 0xEE))
txb(sl, "Perfect: Leather 1.000 | Bottle 0.999 | Metal Nut 0.987    Only screw (0.717) & transistor (0.736) remain challenging",
    7.2, 7.03, 5.6, 0.32, size=9, color=BLACK, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — THREE-WAY COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Three-Way Comparison — Autoencoder vs PatchCore vs I-JEPA + ViT")

# Top: two images side by side
if IMGS["3way_compare"]:
    add_img(sl, "3way_compare", 0.35, 1.28, 6.3, 3.5)
if IMGS["mean_auroc_bar"]:
    add_img(sl, "mean_auroc_bar", 6.85, 1.28, 6.1, 3.5)

# Bottom: summary table
rect(sl, 0.35, 4.92, 12.6, 0.48, ACCENT)
for ci, (label, cx, cw) in enumerate([
    ("Method", 0.45, 2.8), ("Mean Image AUROC", 3.4, 2.3), ("Mean Pixel AUROC", 5.85, 2.3),
    ("Training Needed", 8.25, 2.1), ("Core Signal", 10.5, 2.35)]):
    txb(sl, label, cx, 4.97, cw, 0.38, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

comp_rows = [
    ("Autoencoder (Phase 2)",       "0.659", "0.737", "Yes — 50 epochs",   "Pixel MSE",      RGBColor(0xFF,0xF3,0xF3), RED),
    ("PatchCore (ResNet backbone)", "0.920", "0.870", "No — pretrained",   "Patch k-NN",     RGBColor(0xF5,0xF5,0xF5), GRAY),
    ("I-JEPA + ViT (Ours)  ★",     "0.919", "0.892", "No — pretrained",   "Semantic k-NN",  RGBColor(0xF0,0xFF,0xF4), GREEN),
]
for ri, (method, img_a, pix_a, train, signal, bg, ac) in enumerate(comp_rows):
    ty = 5.44 + ri * 0.52
    rect(sl, 0.35, ty, 12.6, 0.50, bg)
    bold_r = ri == 2
    txb(sl, method, 0.45, ty+0.1, 2.9, 0.35, size=10, bold=bold_r, color=ac)
    txb(sl, img_a,  3.4,  ty+0.1, 2.3, 0.35, size=11, bold=bold_r, color=ac, align=PP_ALIGN.CENTER)
    txb(sl, pix_a,  5.85, ty+0.1, 2.3, 0.35, size=11, bold=bold_r, color=ac, align=PP_ALIGN.CENTER)
    txb(sl, train,  8.25, ty+0.1, 2.1, 0.35, size=10, bold=bold_r, color=ac, align=PP_ALIGN.CENTER)
    txb(sl, signal, 10.5, ty+0.1, 2.35,0.35, size=10, bold=bold_r, color=ac, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONSISTENCY & VARIANCE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Why Consistency Matters — Reliability Across All 15 Categories")

if IMGS["boxplot"]:
    add_img(sl, "boxplot", 0.35, 1.28, 5.8, 4.5)
if IMGS["pixel_var"]:
    add_img(sl, "pixel_var", 6.5, 1.28, 6.4, 4.5)

# Bottom observations
rect(sl, 0.35, 5.88, 5.8, 1.3, RGBColor(0xFF, 0xF3, 0xF3))
txb(sl, "Autoencoder Distribution", 0.55, 5.95, 5.4, 0.38, size=12, bold=True, color=RED)
txb_lines(sl, [
    ("Median ~0.65   |   Range: 0.37 – 0.98", False, BLACK, 11),
    ("Huge spread → completely unreliable for production", True, RED, 11),
], 0.55, 6.36, 5.5, 0.75, size=11)

rect(sl, 6.5, 5.88, 6.4, 1.3, RGBColor(0xF0, 0xFF, 0xF4))
txb(sl, "ViT-PatchCore: Tightest Distribution", 6.7, 5.95, 6.1, 0.38, size=12, bold=True, color=GREEN)
txb_lines(sl, [
    ("Pixel AUROC Variance: 0.0025   (vs AE: 0.027  =  10× lower)", True, GREEN, 11),
    ("Median ~0.95 with very tight IQR → reliable on ALL categories", False, BLACK, 11),
], 6.7, 6.36, 6.1, 0.75, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — IMPROVEMENT ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Improvement Analysis — ViT vs Autoencoder Per Category")

if IMGS["improvement"]:
    add_img(sl, "improvement", 0.35, 1.28, 7.5, 5.35)

rect(sl, 8.2, 1.28, 4.75, 5.35, RGBColor(0xF5, 0xF8, 0xFF))
txb(sl, "Top Improvements", 8.4, 1.36, 4.45, 0.42, size=13, bold=True, color=ACCENT)
top5 = [
    ("cable",     "0.37", "0.94", "+0.57"),
    ("leather",   "0.84", "1.00", "+0.55"),
    ("carpet",    "0.43", "0.97", "+0.54"),
    ("bottle",    "0.67", "1.00", "+0.33"),
    ("metal nut", "0.44", "0.99", "+0.27"),
]
rect(sl, 8.2, 1.82, 4.75, 0.38, ACCENT)
for ci, (lbl, cx) in enumerate([("Category", 8.3), ("AE", 9.5), ("ViT", 10.3), ("Gain", 11.1)]):
    txb(sl, lbl, cx, 1.87, 0.85, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, (cat, ae, vit, gain) in enumerate(top5):
    ty = 2.24 + ri * 0.44
    rect(sl, 8.2, ty, 4.75, 0.42, RGBColor(0xF0,0xFF,0xF4) if ri%2==0 else WHITE)
    txb(sl, cat,  8.3,  ty+0.08, 1.15, 0.28, size=10, color=BLACK)
    txb(sl, ae,   9.5,  ty+0.08, 0.75, 0.28, size=10, color=RED, align=PP_ALIGN.CENTER)
    txb(sl, vit,  10.3, ty+0.08, 0.75, 0.28, size=10, color=GREEN, align=PP_ALIGN.CENTER)
    txb(sl, gain, 11.1, ty+0.08, 0.78, 0.28, size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txb_lines(sl, [
    ("Minimal gain categories:", True, BLACK, 11),
    ("• Hazelnut: AE=0.98 → ViT=0.98 (+0.0)", False, GRAY, 10),
    ("  (AE already near-perfect here)", False, GRAY, 10),
    ("• Wood: AE=0.93 → ViT=0.98 (+0.05)", False, GRAY, 10),
    ("", False, GRAY, 4),
    ("Key insight:", True, BLACK, 11),
    ("ViT makes weak categories strong", False, GRAY, 10),
    ("without degrading strong ones.", False, GRAY, 10),
], 8.4, 4.48, 4.45, 2.0, size=10)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — TREND + SCATTER
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Trend & Scatter Analysis — Stability and Spatial Localization")

if IMGS["trend"]:
    add_img(sl, "trend", 0.35, 1.28, 6.3, 4.5)
if IMGS["scatter"]:
    add_img(sl, "scatter", 6.85, 1.28, 6.1, 4.5)

rect(sl, 0.35, 5.88, 6.3, 1.3, RGBColor(0xF5, 0xF8, 0xFF))
txb(sl, "Trend Across 15 Categories", 0.55, 5.95, 6.0, 0.38, size=11, bold=True, color=ACCENT)
txb_lines(sl, [
    ("AE (blue): wild swings (0.37–0.98).   PatchCore & ViT (orange/green): flat at top.", False, BLACK, 10),
    ("Flatness = reliability. Screw is the hardest for all methods (fine-grained threads).", False, GRAY, 10),
], 0.55, 6.36, 6.0, 0.75, size=10)

rect(sl, 6.85, 5.88, 6.1, 1.3, RGBColor(0xF0, 0xFF, 0xF4))
txb(sl, "Image vs Pixel AUROC — Detect AND Localize", 7.05, 5.95, 5.8, 0.38, size=11, bold=True, color=GREEN)
txb_lines(sl, [
    ("Top-right = detect AND localize correctly.   ViT dots cluster top-right.", False, BLACK, 10),
    ("Our method doesn't just flag anomalies — it shows WHERE the defect is.", True, GREEN, 10),
], 7.05, 6.36, 5.8, 0.75, size=10)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — ALL MODELS COMPARISON HEATMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Complete Results Table — All Methods × All Categories × Both Metrics")

if IMGS["all_models_hm"]:
    add_img(sl, "all_models_hm", 0.35, 1.3, 8.2, 5.7)

rect(sl, 8.85, 1.3, 4.1, 5.7, RGBColor(0xF5, 0xF8, 0xFF))
txb(sl, "How to Read", 9.05, 1.38, 3.8, 0.42, size=13, bold=True, color=ACCENT)
txb_lines(sl, [
    ("Columns:", True, BLACK, 12),
    ("AE = Autoencoder", False, GRAY, 11),
    ("PC = PatchCore (ResNet)", False, GRAY, 11),
    ("VIT = I-JEPA + ViT (Ours)", False, GRAY, 11),
    ("_Image = detection AUROC", False, GRAY, 11),
    ("_Pixel = localization AUROC", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("Colours:", True, BLACK, 12),
    ("Dark red = near 1.0 (excellent)", False, RGBColor(0xC0,0x39,0x2B), 11),
    ("Blue = lower scores (poor)", False, TEAL, 11),
    ("", False, GRAY, 5),
    ("Key takeaway:", True, BLACK, 12),
    ("AE column is mostly pale/blue.", False, GRAY, 11),
    ("VIT column is mostly dark red.", False, GREEN, 11),
    ("", False, GRAY, 5),
    ("VIT consistently outperforms", True, ACCENT, 11),
    ("AE across nearly all categories", True, ACCENT, 11),
    ("for both metrics.", True, ACCENT, 11),
], 9.05, 1.85, 3.8, 5.0, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — RADAR CHART
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Radar Chart — ViT vs PatchCore Across All 15 Categories")

if IMGS["radar"]:
    add_img(sl, "radar", 1.5, 1.2, 10.2, 5.5)

rect(sl, 0.35, 6.82, 12.6, 0.55, RGBColor(0xF0, 0xF7, 0xFF))
txb(sl, "ViT and PatchCore fill almost the entire radar — both near 1.0 on most axes.  Only screw, transistor, and grid show minor dips.  AE (not shown) would fill only ~65% of the area.",
    0.5, 6.87, 12.2, 0.45, size=10, color=DARK, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — AUTOENCODER HEATMAPS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Heatmaps — Autoencoder: Pixel Reconstruction Error (Leather Category)")

hm_ae = [
    ("hm_cut",     "Leather — Cut defect",      "Diffuse error across non-defect\nareas too — noisy localization"),
    ("hm_fold",    "Leather — Fold defect",      "Error spreads beyond the fold.\nFalse positives in background"),
    ("hm_glue",    "Leather — Glue defect",      "Some signal, but background\nnoise contaminates the heatmap"),
    ("hm_good_ae", "Leather — NORMAL image",     "Should be all-blue.\nFew false-alarm patches visible"),
]
for i, (key, title, note) in enumerate(hm_ae):
    xl = 0.35 + i * 3.25
    if IMGS.get(key):
        add_img(sl, key, xl, 1.3, 3.0, 3.7)
    else:
        rect(sl, xl, 1.3, 3.0, 3.7, RGBColor(0xEE,0xEE,0xEE))
    txb(sl, title, xl, 5.08, 3.0, 0.38, size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txb(sl, note, xl, 5.48, 3.0, 0.55, size=9, color=GRAY, align=PP_ALIGN.CENTER)

rect(sl, 0.35, 6.1, 12.6, 0.98, RGBColor(0xFF, 0xF3, 0xF3))
txb(sl, "Observation:", 0.55, 6.16, 2.0, 0.35, size=11, bold=True, color=RED)
txb_lines(sl, [
    ("Pixel MSE heatmaps are noisy — reconstruction error spreads across the whole image, not just the defect region.", False, BLACK, 11),
    ("The model reconstructs defects somewhat well, so the error signal is weak and un-localized.  Qualitative localization is poor.", False, GRAY, 11),
], 2.6, 6.16, 10.2, 0.88, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — I-JEPA + VIT HEATMAPS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Heatmaps — I-JEPA + ViT: Semantic Patch-Level Localization")

hm_vit = [
    ("hm_ijepa_hz1", "Hazelnut — Crack",    "Hot patches align\nprecisely with crack"),
    ("hm_ijepa_hz2", "Hazelnut — Defect",   "Clean separation:\nred=defect, blue=normal"),
    ("hm_ijepa_cb1", "Cable — Defect",       "Wire anomaly region\ncorrectly highlighted"),
    ("hm_ijepa_cb2", "Cable — Defect 2",     "Fine-grained cable\ndefect localized"),
    ("hm_ijepa_mn",  "Metal Nut — Defect",   "Scratch/pit isolated\nto correct patch"),
    ("hm_ijepa_sc",  "Screw — Defect",       "Hardest category;\npartial localization"),
]
for i, (key, title, note) in enumerate(hm_vit):
    col = i % 3
    row = i // 3
    xl = 0.35 + col * 4.3
    yt = 1.28 + row * 3.25
    if IMGS.get(key):
        add_img(sl, key, xl, yt, 4.0, 2.7)
    else:
        rect(sl, xl, yt, 4.0, 2.7, RGBColor(0xEE,0xEE,0xEE))
    txb(sl, title, xl, yt+2.72, 4.0, 0.32, size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txb(sl, note,  xl, yt+3.02, 4.0, 0.38, size=9,  color=GRAY, align=PP_ALIGN.CENTER)

rect(sl, 0.35, 7.02, 12.6, 0.42, RGBColor(0xF0, 0xFF, 0xF4))
txb(sl, "Each 16×16 pixel patch is independently scored against the memory bank.  Red = high anomaly score, Blue = low (normal).  No pixel reconstruction involved — pure semantic distance.",
    0.5, 7.05, 12.2, 0.35, size=9, color=DARK, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — WEB APPLICATION
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Real-Time Web Application — Live Defect Detection Demo")

# Left
rect(sl, 0.4, 1.28, 5.8, 5.85, RGBColor(0xF5, 0xF8, 0xFF))
txb(sl, "Application Architecture", 0.6, 1.36, 5.5, 0.42, size=13, bold=True, color=ACCENT)
txb_lines(sl, [
    ("Backend:", True, BLACK, 12),
    ("• Python Flask server  (app.py + backend.py)", False, GRAY, 11),
    ("• Auto model selection:", False, GRAY, 11),
    ("  If ViT detector .pkl exists → uses I-JEPA+ViT", False, GREEN, 11),
    ("  Otherwise → falls back to Autoencoder", False, GRAY, 11),
    ("• Supports all 15 MVTec categories", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("Frontend:", True, BLACK, 12),
    ("• Pure HTML/CSS/JavaScript — no framework", False, GRAY, 11),
    ("• Drag-and-drop image upload", False, GRAY, 11),
    ("• Real-time anomaly heatmap display", False, GRAY, 11),
    ("", False, GRAY, 5),
    ("How to run:", True, BLACK, 12),
    ("python app.py", False, ACCENT, 12),
    ("Open: http://127.0.0.1:5000", False, ACCENT, 11),
    ("", False, GRAY, 5),
    ("Output per image:", True, BLACK, 12),
    ("• Anomaly score (0–1)", False, GRAY, 11),
    ("• Normal / Anomaly badge", False, GRAY, 11),
    ("• Heatmap overlay (saved as PNG)", False, GRAY, 11),
    ("• ROC-AUC and F1 for the category", False, GRAY, 11),
], 0.6, 1.82, 5.5, 5.2, size=11)

# Right: two tabs
rect(sl, 6.7, 1.28, 6.2, 2.75, RGBColor(0xE8, 0xF5, 0xFF))
txb(sl, "Tab 1 — Image Tester", 6.9, 1.36, 5.9, 0.42, size=12, bold=True, color=TEAL)
txb_lines(sl, [
    ("• Drag & drop any image (or browse files)", False, GRAY, 11),
    ("• Select product category from dropdown", False, GRAY, 11),
    ("• Click 'Run Detection'", False, GRAY, 11),
    ("• See: score, Normal/Anomaly badge, heatmap,", False, GRAY, 11),
    ("  ROC-AUC, F1, metric explanations", False, GRAY, 11),
    ("• Ideal for testing custom/real-world images", False, GRAY, 11),
], 6.9, 1.82, 5.9, 2.1, size=11)

rect(sl, 6.7, 4.22, 6.2, 2.91, RGBColor(0xF0, 0xFF, 0xF0))
txb(sl, "Tab 2 — Dataset Browser", 6.9, 4.3, 5.9, 0.42, size=12, bold=True, color=GREEN)
txb_lines(sl, [
    ("• Browse all 1,725 MVTec test images", False, GRAY, 11),
    ("• Filter by category (any of 15)", False, GRAY, 11),
    ("• Click any image → instant detection", False, GRAY, 11),
    ("• Perfect for live demonstration:", False, GRAY, 11),
    ("  show good vs defective side-by-side", False, GRAY, 11),
    ("• Heatmaps saved to results/heatmaps/", False, GRAY, 11),
], 6.9, 4.76, 5.9, 2.28, size=11)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — KEY FINDINGS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Key Findings & Discussion")

findings = [
    ("1", "Pixel MSE is Fundamentally Limited",
     "Mean AE AUROC = 0.659 with extreme variance (0.37–0.98). Autoencoders generalize to reconstruct anomalies well → reconstruction error is an unreliable anomaly signal. Confirms literature findings.",
     RED),
    ("2", "Semantic Features Unlock Accuracy",
     "I-JEPA principle: compare patch embeddings, not pixels. Moving to feature space = +26% mean AUROC. Normal patches cluster tightly in ViT space; anomaly patches are clear outliers.",
     ACCENT),
    ("3", "Pretrained Backbone > Scratch Training for Small Datasets",
     "300–400 normal images per category is insufficient for self-supervised ViT training from scratch. ImageNet-pretrained ViT features generalize immediately to industrial inspection. Transfer learning wins.",
     TEAL),
    ("4", "Patch-Level Scoring = Detection + Localization",
     "Not just 'is this image anomalous?' but 'which 16×16 pixel region is the defect?' Pixel AUROC = 0.892 means the system accurately localizes defects — directly useful for operators.",
     GREEN),
    ("5", "Variance Matters as Much as Mean",
     "ViT Pixel AUROC variance = 0.0025 vs AE's 0.027 (10× lower). A reliable production system must work on ALL product types, not just easy ones. Consistency is a key engineering metric.",
     ORANGE),
]
for i, (num, title, body, c) in enumerate(findings):
    ty = 1.28 + i * 1.18
    rect(sl, 0.35, ty, 0.62, 1.08, c)
    txb(sl, num, 0.35, ty+0.28, 0.62, 0.52, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, 0.98, ty, 12.0, 1.08, RGBColor(0xF9,0xF9,0xFF) if i%2==0 else WHITE)
    txb(sl, title, 1.12, ty+0.08, 11.7, 0.38, size=13, bold=True, color=c)
    txb(sl, body, 1.12, ty+0.46, 11.7, 0.55, size=10, color=GRAY)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Conclusion — What This Project Achieved")

cols_c = [
    ("The Goal", ACCENT, [
        "Build an AI system that detects",
        "manufacturing defects without",
        "ever seeing a defective image",
        "during training.",
        "",
        "Unsupervised — labels used only",
        "for evaluation, never for training.",
        "",
        "Real-time — web application for",
        "live inspection with heatmap",
        "localization output.",
    ]),
    ("What We Built", TEAL, [
        "Phase 0 & 1: Environment + MVTec",
        "data pipeline (15 categories)",
        "",
        "Phase 2: Autoencoder baseline",
        "(pixel MSE, benchmark established)",
        "",
        "Phase 3: Custom I-JEPA ViT-Small",
        "(learned the methodology)",
        "",
        "Phase 4: I-JEPA + Pretrained ViT",
        "(production-grade accuracy)",
        "",
        "Web UI: real-time Flask app",
    ]),
    ("What We Achieved", GREEN, [
        "Mean Image AUROC:  0.919",
        "Mean Pixel AUROC:  0.892",
        "",
        "Perfect: Leather  1.000",
        "Perfect: Bottle   0.999",
        "Perfect: Metal nut 0.987",
        "",
        "+26% over Autoencoder baseline",
        "10× lower variance than AE",
        "",
        "All 15 categories evaluated.",
        "Real-time web demo working.",
        "Heatmap localization included.",
    ]),
]
for i, (title, c, items) in enumerate(cols_c):
    xl = 0.35 + i * 4.33
    rect(sl, xl, 1.28, 4.2, 5.85, c)
    txb(sl, title, xl+0.1, 1.35, 4.0, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    lns = [(it, False, WHITE, 11) for it in items]
    txb_lines(sl, lns, xl+0.2, 1.88, 3.9, 5.1, size=11)

# Big bottom summary
rect(sl, 0.35, 7.06, 12.6, 0.4, RGBColor(0x0F,0x3D,0x91))
txb(sl,
    "We built an industry-grade AI system that detects and localizes manufacturing defects in real-time with 91.9% accuracy — using zero defect images during training.",
    0.5, 7.09, 12.2, 0.33, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "Future Work — What Comes Next")

future = [
    ("1", "Train I-JEPA on Larger Industrial Data",
     "The custom I-JEPA architecture (Phase 3) would perform at par with pretrained ViT if trained on a large-scale industrial dataset. Continue self-supervised pre-training on multi-domain manufacturing images."),
    ("2", "Multi-Layer Feature Fusion",
     "Currently using only the final ViT-B/16 layer. Averaging features from the last 3 transformer blocks (already implemented: extract_multilayer_features()) captures both fine-grained (early) and semantic (late) features simultaneously."),
    ("3", "Video / Sequential Inspection",
     "Extend to video streams for conveyor-belt inspection in real manufacturing lines. Temporal consistency across frames can reduce false positives further."),
    ("4", "Custom Industrial Dataset Integration",
     "Apply the pipeline to real manufacturing datasets from production partners — PCBs, semiconductor wafers, pharmaceutical packaging — beyond MVTec AD."),
    ("5", "Model Compression for Edge Deployment",
     "ViT-B/16 is 86M parameters. For embedded inspection cameras, distill into a MobileViT or EfficientViT while preserving anomaly detection accuracy."),
    ("6", "Few-Shot Defect Fine-Tuning",
     "Allow 3–5 labelled defect images to fine-tune the anomaly threshold per defect type — bridging unsupervised and weakly-supervised for critical applications."),
]
for i, (num, title, body) in enumerate(future):
    col = i % 2
    row = i // 2
    xl = 0.35 + col * 6.5
    ty = 1.28 + row * 2.0
    c = [ACCENT, TEAL, GREEN, ORANGE, RGBColor(0x80,0x40,0x80), RGBColor(0xC0,0x39,0x2B)][i]
    rect(sl, xl, ty, 0.55, 1.85, c)
    txb(sl, num, xl, ty+0.55, 0.55, 0.52, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, xl+0.55, ty, 5.72, 1.85, RGBColor(0xF8,0xF8,0xFF) if i%2==0 else WHITE)
    txb(sl, title, xl+0.7, ty+0.12, 5.5, 0.42, size=12, bold=True, color=c)
    txb(sl, body, xl+0.7, ty+0.58, 5.5, 1.18, size=10, color=GRAY)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — REFERENCES
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
slide_header(sl, "References")

refs = [
    ("1", "I-JEPA", ACCENT,
     "Assran, M., Duval, Q., Misra, I., Bojanowski, P., Vincent, P., Rabbat, M., & LeCun, Y. (2023). Self-Supervised Learning from Images with a Joint-Embedding Predictive Architecture. Proceedings of the IEEE/CVF Conference on Computer Vision and Pattern Recognition (CVPR), 2023. arXiv:2301.08243"),
    ("2", "MVTec AD Dataset", TEAL,
     "Bergmann, P., Batzner, K., Fauser, M., Sattlegger, D., & Steger, C. (2021). The MVTec Anomaly Detection Dataset: A Comprehensive Real-World Dataset for Unsupervised Anomaly Detection. International Journal of Computer Vision, 129(4), 1038–1059. https://doi.org/10.1007/s11263-020-01400-4"),
    ("3", "PatchCore", GREEN,
     "Roth, K., Pemula, L., Zepeda, J., Schölkopf, B., Brox, T., & Gehler, P. (2022). Towards Total Recall in Industrial Anomaly Detection. Proceedings of the IEEE/CVF Conference on Computer Vision and Pattern Recognition (CVPR), 2022."),
    ("4", "Vision Transformer (ViT)", ORANGE,
     "Dosovitskiy, A., Beyer, L., Kolesnikov, A., Weissenborn, D., Zhai, X., Unterthiner, T., ... & Houlsby, N. (2021). An Image is Worth 16×16 Words: Transformers for Image Recognition at Scale. International Conference on Learning Representations (ICLR), 2021."),
    ("5", "Unreliability of Autoencoders", RED,
     "Baur, C., Wiestler, B., Albarqouni, S., & Navab, N. (2021). Autoencoders for Unsupervised Anomaly Segmentation in Brain MR Images: A Comparative Study. Medical Image Analysis, 69, 101952. (The statistical evidence supporting our motivation to move beyond pixel reconstruction.)"),
],
for reflist in refs:
    for i, (num, label, c, body) in enumerate(reflist):
        ty = 1.28 + i * 1.18
        rect(sl, 0.35, ty, 0.5, 1.08, c)
        txb(sl, num, 0.35, ty+0.28, 0.5, 0.52, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, 0.86, ty, 12.1, 1.08, RGBColor(0xF7,0xF9,0xFF) if i%2==0 else WHITE)
        txb(sl, label, 1.0, ty+0.08, 11.8, 0.36, size=12, bold=True, color=c)
        txb(sl, body,  1.0, ty+0.46, 11.8, 0.58, size=9,  color=GRAY)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = os.path.join(BASE, "Final_Presentation_GroupA5.pptx")
prs.save(out)
print(f"✓  Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
